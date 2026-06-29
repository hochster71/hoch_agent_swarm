import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from backend.artifacts.brand_renderer import BrandRenderer

class SlideFactory:
    def __init__(self, root_dir="/Users/michaelhoch/hoch_agent_swarm"):
        self.root_dir = root_dir
        self.renderer = BrandRenderer(root_dir)

    def create_deck(self, title: str, subtitle: str, slides_content: list, filepath: str):
        prs = Presentation()
        
        # Slide Dimensions
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625) # 16:9 widescreen
        
        # Colors from brand
        colors = self.renderer.get_theme_colors()
        # Parse hex colors to RGB
        bg_rgb = self._hex_to_rgb(colors.get("background", "#0a0a0c"))
        teal_rgb = self._hex_to_rgb(colors.get("accent_teal", "#10b981"))
        text_rgb = self._hex_to_rgb(colors.get("text_primary", "#f3f4f6"))
        
        # 1. Title Slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Set dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_rgb)
        
        # Title text box
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(*teal_rgb)
        p.font.name = "Outfit"
        
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(*text_rgb)
        p2.font.name = "Inter"
        
        # 2. Content Slides
        for s_title, bullet_points in slides_content:
            slide = prs.slides.add_slide(blank_slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_rgb)
            
            # Content title
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = s_title
            p_title.font.bold = True
            p_title.font.size = Pt(28)
            p_title.font.color.rgb = RGBColor(*teal_rgb)
            p_title.font.name = "Outfit"
            
            # Bullets
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(3.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            for pt in bullet_points:
                p_bullet = tf_content.add_paragraph()
                p_bullet.text = f"• {pt}"
                p_bullet.font.size = Pt(16)
                p_bullet.font.color.rgb = RGBColor(*text_rgb)
                p_bullet.font.name = "Inter"
                
        # Ensure output folder exists
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        prs.save(filepath)
        return filepath

    def _hex_to_rgb(self, hex_str: str) -> tuple:
        hex_str = hex_str.lstrip('#')
        if len(hex_str) == 3:
            hex_str = ''.join([c*2 for c in hex_str])
        return int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
